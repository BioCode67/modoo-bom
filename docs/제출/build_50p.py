# -*- coding: utf-8 -*-
"""모두봄 — 최종제출용 50p PPTX 생성기.
검증된 사실·실측 수치만 사용(과장 0). 이미지: 기존 차트/다이어그램/실브라우저 캡처.
실행: python docs/제출/build_50p.py  → docs/제출/모두봄-최종제출-50p.pptx
"""
import sys, os
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
A = lambda *p: os.path.join(ROOT, *p)  # 리포지토리 상대 경로

# ── 팔레트 ──────────────────────────────────────────────
GREEN_D = RGBColor(0x15, 0x80, 0x3D)   # sprout-700
GREEN = RGBColor(0x22, 0xC5, 0x5E)     # sprout-500
CREAM = RGBColor(0xFF, 0xFA, 0xF3)
INK = RGBColor(0x1A, 0x25, 0x1E)
GRAY = RGBColor(0x55, 0x5F, 0x58)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PEACH = RGBColor(0xC2, 0x41, 0x0C)

W, H = Inches(13.333), Inches(7.5)
FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
page_no = 0

def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def _text(tf, runs, size, color, bold=False, align=PP_ALIGN.LEFT, space_after=6):
    """runs: str 또는 [(text, bold?, color?)] 리스트의 리스트(문단별)"""
    tf.word_wrap = True
    if isinstance(runs, str):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        if isinstance(para, str):
            para = [(para, bold, color)]
        for seg in para:
            t, b, c = (seg + (None,) * 3)[:3] if isinstance(seg, tuple) else (seg, bold, color)
            r = p.add_run(); r.text = t
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.bold = b if b is not None else bold
            r.font.color.rgb = c if c is not None else color

def base_slide(title=None, eyebrow=None):
    """공통 골격: 크림 배경 + 상단 타이틀 바 + 하단 푸터/쪽번호"""
    global page_no
    page_no += 1
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, W, H); _fill(bg, CREAM)
    if title:
        bar = s.shapes.add_shape(1, 0, 0, W, Inches(1.0)); _fill(bar, GREEN_D)
        tb = s.shapes.add_textbox(Inches(0.55), Inches(0.08), Inches(12.2), Inches(0.86))
        tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        if eyebrow:
            _text(tf, [[(eyebrow + "  ", False, RGBColor(0xBB, 0xF7, 0xD0))], [(title, True, WHITE)]], 13, WHITE)
            tf.paragraphs[1].runs[0].font.size = Pt(24)
        else:
            _text(tf, title, 24, WHITE, bold=True)
    # 푸터
    ft = s.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.35))
    _text(ft.text_frame, [[("모두봄 — 개인 복지 자산 관리 AI Agent", False, GRAY), (f"        {page_no:02d} / 50", False, GRAY)]], 10, GRAY)
    return s

def bullets(s, items, x=0.7, y=1.35, w=11.9, h=5.5, size=16, gap=10):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        if isinstance(it, tuple):  # (text, indent?, bold?, color?)
            text = it[0]
            lvl = it[1] if len(it) > 1 else 0
            bold = it[2] if len(it) > 2 else False
            col = it[3] if len(it) > 3 else INK
        else:
            text, lvl, bold, col = it, 0, False, INK
        p.level = lvl
        r = p.add_run(); r.text = ("· " if lvl == 0 else "– ") + text if not text.startswith(("✓", "⚠", "★", "①", "②", "③", "④", "⑤")) else text
        r.font.name = FONT; r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.bold = bold; r.font.color.rgb = col
    return tb

def img_slide(title, img, caption=None, eyebrow=None, note=None):
    s = base_slide(title, eyebrow)
    if os.path.exists(img):
        from PIL import Image
        iw, ih = Image.open(img).size
        maxw, maxh = Inches(11.6), Inches(5.05 if caption or note else 5.5)
        ratio = min(maxw / Emu(1) / iw, maxh / Emu(1) / ih)
        pw, ph = int(iw * ratio), int(ih * ratio)
        left = int((W - pw) / 2)
        s.shapes.add_picture(img, left, Inches(1.25), width=pw, height=ph)
        ytxt = Inches(1.25) + ph + Inches(0.08)
    else:
        tb = s.shapes.add_textbox(Inches(0.7), Inches(3), Inches(11.9), Inches(1))
        _text(tb.text_frame, f"[이미지 준비: {os.path.basename(img)}]", 14, GRAY)
        ytxt = Inches(4.2)
    if caption:
        cb = s.shapes.add_textbox(Inches(0.7), ytxt, Inches(11.9), Inches(0.5))
        _text(cb.text_frame, caption, 13, GRAY, align=PP_ALIGN.CENTER)
    if note:
        nb = s.shapes.add_textbox(Inches(0.7), ytxt + Inches(0.42), Inches(11.9), Inches(0.6))
        _text(nb.text_frame, note, 12, GREEN_D, align=PP_ALIGN.CENTER, bold=True)
    return s

def stat_cards(s, cards, y=1.6, ch=1.6, size_v=30, size_l=13):
    n = len(cards); gapx = Inches(0.25)
    cw = int((W - Inches(1.1) - gapx * (n - 1)) / n)
    x = Inches(0.55)
    for v, l in cards:
        box = s.shapes.add_shape(1, x, Inches(y), cw, Inches(ch))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.fill.solid(); box.line.fill.fore_color.rgb = GREEN; box.line.width = Pt(1.2)
        tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(tf, [[(v, True, GREEN_D)]], size_v, GREEN_D, align=PP_ALIGN.CENTER, space_after=2)
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = l; r.font.name = FONT; r.font.size = Pt(size_l); r.font.color.rgb = GRAY
        x += cw + gapx

def section(no, title, sub):
    """섹션 구분 슬라이드"""
    global page_no
    page_no += 1
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, W, H); _fill(bg, GREEN_D)
    tb = s.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.4))
    tf = tb.text_frame
    _text(tf, [[(f"PART {no}", True, RGBColor(0xBB, 0xF7, 0xD0))]], 18, WHITE, space_after=10)
    p = tf.add_paragraph(); r = p.add_run(); r.text = title
    r.font.name = FONT; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    r2 = p2.add_run(); r2.text = sub
    r2.font.name = FONT; r2.font.size = Pt(16); r2.font.color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
    ft = s.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.35))
    _text(ft.text_frame, f"모두봄        {page_no:02d} / 50", 10, RGBColor(0xD1, 0xFA, 0xE5))
    return s

CAP = lambda *p: A("frontend", "e2e", "audit", *p)
DOC = lambda *p: A("docs", "기획서자료", *p)
SUB = lambda *p: A("docs", "제출", "스크린샷", *p)

# ════════════════════════════════ 1. 표지 ════════════════════════════════
s = prs.slides.add_slide(BLANK); page_no += 1
bg = s.shapes.add_shape(1, 0, 0, W, H); _fill(bg, GREEN_D)
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.6))
tf = tb.text_frame
_text(tf, [[("2026 AI·SW중심대학 디지털 경진대회 · SW부문 최종 산출물", False, RGBColor(0xBB, 0xF7, 0xD0))]], 16, WHITE, space_after=16)
p = tf.add_paragraph(); r = p.add_run(); r.text = "모두봄 🌱"
r.font.name = FONT; r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = WHITE
p = tf.add_paragraph(); p.space_before = Pt(6)
r = p.add_run(); r.text = "몰라서 못 받는 복지를 찾아 신청까지 — 개인 복지 자산 관리 AI Agent"
r.font.name = FONT; r.font.size = Pt(22); r.font.color.rgb = WHITE
p = tf.add_paragraph(); p.space_before = Pt(24)
r = p.add_run(); r.text = "라이브: biocode67.github.io/modoo-bom   ·   설치·로그인 없이 지금 바로 체험"
r.font.name = FONT; r.font.size = Pt(15); r.font.color.rgb = RGBColor(0xD1, 0xFA, 0xE5)

# ════════════════════════ 2. Summary ════════════════════════
s = base_slide("Summary — 한 장 요약")
bullets(s, [
    ("문제: 자격이 있어도 몰라서·어려워서·못 써서 놓치는 복지 — 정부·지자체·민간에 5,300여 건이 흩어져 있습니다.", 0, True),
    ("해법: 말 한마디(음성·7개 언어)면 AI Agent가 찾고(인지·판단), 담고·서류 만들고·발급/신청까지 실행(행동), 마감·갱신을 감시(관찰)합니다.", 0),
    ("차별점 ① 정부 데이터 밖 민간재단 21건 실측 큐레이션  ② 온디바이스 AI(개인정보 서버 전송 0)  ③ 실제 정부24를 조작하는 서류 15종 자동발급  ④ 창구 도움 카드·사유서 대필 — '손에 쥐여주는' 마지막 1cm", 0),
    ("정직성: 과장을 코드로 차단(현금성 게이트·심사형 '자격 있음' 표기 거부) — 다중에이전트 적대 감사 17라운드로 확정 결함 90여 건 정정.", 0),
])
stat_cards(s, [("5,300+", "실데이터 복지(건)"), ("809", "자동 테스트(프론트 679·백엔드 130)"), ("15종", "서류 자동발급"), ("7개 언어", "온디바이스 AI 검색")], y=5.2, ch=1.5, size_v=26)

# ════════════════════════ PART 1 문제 ════════════════════════
section(1, "문제 — 몰라서 못 받는 복지", "복지 사각지대는 '제도 부족'이 아니라 '연결 실패'의 문제입니다")

s = base_slide("자격이 있어도 받지 못하는 3가지 이유")
bullets(s, [
    ("① 몰라서 — 중앙부처 460 + 지자체 4,683 + 민간재단·서민금융까지, 어디에 무엇이 있는지 한 곳에서 볼 수 없습니다.", 0, True),
    ("복지로·정부24·지자체 홈페이지·재단 공고가 제각각 — '내 것'을 찾으려면 수십 개 사이트를 뒤져야 합니다.", 1),
    ("② 어려워서 — '소득인정액', '중위소득 48%', '차상위'… 자격 기준은 행정 용어의 벽 뒤에 있습니다.", 0, True),
    ("③ 못 써서 — 찾아도 신청서의 '사유'를 쓰지 못하고, 창구에서 말이 막히고, 서류 발급에서 포기합니다.", 0, True),
    ("특히 어르신·장애인·외국인·위기가구 — 복지가 가장 필요한 분들이 가장 큰 장벽을 만납니다.", 0, False, PEACH),
])

s = base_slide("한 사람의 이야기 — 김복순 님(72세, 독거)")
bullets(s, [
    ("기초연금은 받고 계십니다. 그런데 그게 전부일까요?", 0, True),
    ("모두봄에 '72세 혼자 사는데 소득이 적어요' 한 문장을 넣으면 —", 0),
    ("핵심 현금지원 월 최대 64만원(중복수급 미반영), 강력 추천 10여 건이 나옵니다: 주거급여·의료급여·틀니 지원·노인 일자리·에너지바우처…", 1, True),
    ("정부 데이터 어디에도 없는 민간재단(심장 수술비·위기가정 지원)까지 함께.", 1),
    ("김복순 님 같은 분이 전국에 수백만 명입니다. 모두봄은 그분들의 '몰랐던 권리'를 찾아드립니다.", 0, True, GREEN_D),
])

# ════════════════════════ PART 2 솔루션 개요 ════════════════════════
section(2, "솔루션 — AI Agent '모두봄'", "찾기 → 이해 → 행동(신청) → 사후관리, 전 여정을 에이전트가 대행합니다")

s = base_slide("모두봄이란 — 챗봇이 아니라 '루프'입니다")
bullets(s, [
    ("개인 복지 자산 관리 AI Agent — 나이·소득·상황만 알려주면 찾고, 계산하고, 서류·신청·사후관리까지 함께합니다. 회원가입·설치·비용 없음(PWA·오프라인 지원).", 0, True),
    ("인지 — 프로필·자연어·음성에서 상황 신호를 읽습니다(실직·한부모·장애·위기 등 사각지대 신호 포함).", 0),
    ("판단 — 2026 정밀 선정기준으로 자격을 판정하고, 우선순위·예상 금액·'왜 나에게 맞는지' 근거를 만듭니다.", 0),
    ("행동 — 한 번에 담기, 서류 자동발급(RPA), 신청서 정보 자동 입력, 사유서·위임장·전화 대본 대필.", 0),
    ("관찰 — 신청 진행·서류 미비·갱신 임박·마감을 감시하고 먼저 보고합니다(챗을 열면 브리핑부터).", 0),
    ("이 루프 전체가 에이전트입니다 — 인증·최종 제출만 사람에게 남깁니다(설계 원칙).", 0, True, PEACH),
])

s = img_slide("전체 아키텍처 — 3층 지능", DOC("diagram1-아키텍처.png"),
    caption="① 온디바이스 신경망(다국어 의미검색·프라이버시) ② LangGraph 10노드 에이전트(클라우드/로컬) ③ Playwright RPA(실제 정부24)")

s = img_slide("LangGraph 10노드 에이전트", DOC("diagram5-LangGraph10노드.png"),
    caption="프로필 분석 → RAG 검색 → 자격 판별 → 재판정 루프(reflection) → 가이드·서류·포트폴리오·알림·추적 → 오케스트레이터",
    note="라이브에서 실시간 스트리밍으로 10노드가 순차 점등 — 연출이 아니라 실제 실행입니다")

s = img_slide("데이터 — 약 5,300건 실데이터", DOC("chart1-출처구성.png"),
    caption="정부 시드 124(정밀 규칙·검증 금액) + 지원사업 33 + 주택 7 + 서민금융 5 + 민간재단 21 + 공공데이터(중앙 460·지자체 4,683)",
    note="가짜 데이터 미생성 원칙 — 전 항목 실데이터, 헤드라인 금액은 2026 공식 출처로 재검증")

s = base_slide("민간재단 큐레이션 21건 — 정부 데이터 밖의 사각지대")
bullets(s, [
    ("복지로·보조금24 어디에도 없는 기업·재단 지원을 하나하나 실측 검증해 수록했습니다.", 0, True),
    ("장학: 현대차 정몽구 스칼러십 · 관정 · 삼성꿈장학 · 미래에셋", 0),
    ("의료·위기: 한국심장재단 · 백혈병어린이재단 · 아산 SOS · 초록우산 · 밀알(장애·재활)", 0),
    ("정직한 표시 — 심사·선발형임을 명시하고 '자격 있음'으로 단정하지 않습니다(테스트가 이를 강제).", 0, False, PEACH),
    ("서민금융 5건(햇살론·소액생계비대출 등)도 '대출'임을 명시하고 현금 합산에서 제외합니다.", 0),
])

# ════════════════════════ PART 3 찾기 ════════════════════════
section(3, "찾기 — 말 한마디면 됩니다", "대화·한 문장·음성·7개 언어, 어떤 방식으로든")

s = img_slide("대화형 온보딩 — '새싹이와 대화'", CAP("m2-analyze-chat.png"),
    caption="폼 대신 마스코트가 하나씩 묻고 탭으로만 답합니다. 어르신(65+)에겐 큰글씨를 원탭 제안, 새로고침해도 이어서.")

s = img_slide("한 문장 분석 — '72세 혼자 사는데 소득이 적어요'", SUB("한문장분석결과-2026-07-15.png"),
    caption="핵심 현금지원 월 최대 64만원 — 2026-07-15 라이브 실측 화면 그대로(수치·문구 무보정)")

s = img_slide("분석 결과 — 에이전트가 '다음 행동'을 제안", CAP("10-result-top.png"),
    caption="'이렇게 이해했어요' 신호 투명화 + 추천 한 번에 담기 + 마감 임박·놓친 복지·연계 감면까지 이어집니다")

s = img_slide("정책 상세 — '내가 받을 수 있는 이유'", CAP("d1-detail-top.png"),
    caption="조건 나열이 아니라 내 정보와 대조한 체크리스트(만 72세 ✓ · 소득 하위 25% ✓)로 근거를 보여줍니다")

s = base_slide("숨은 지능 — 추천에서 끝나지 않는 판단")
bullets(s, [
    ("수급 조합 도우미 — '하나만'(기초연금↔장애인연금) · '감액 주의'(기초연금→생계급여) · '병급 가능'(아동수당+부모급여). 공식 규정 실측.", 0),
    ("아깝게 놓친 복지 — 딱 한 조건만 벗어난 정책을 근거와 함께(판단 근거 노출).", 0),
    ("연계 감면 일괄 안내 — 자격이 되면 딸려오는 통신·전기·가스·TV수신료 감면을 묶어 안내.", 0),
    ("생애 타임라인·시뮬레이터 — '만 65세가 되면', '아이가 태어나면' 무엇이 열리고 닫히는지 미리.", 0),
    ("긴급복지 빠른 진단 — 위기 상황이면 129·1366·1393 핫라인을 최우선으로(안전 우선 설계).", 0, False, PEACH),
])

s = img_slide("행동형 챗 에이전트 — '다 담아줘' 한마디", CAP("32-chat-result.png"),
    caption="대화 맥락을 기억해 방금 보여준 복지를 저장하고, 서류 요약·자격 설명까지 — 행동은 기기 안 에이전트가 즉시")

s = base_slide("챗 하이브리드 — 행동은 로컬, 지식은 LLM")
bullets(s, [
    ("행동 의도(담기·서류·자격·브리핑) — 기기 안 규칙 에이전트가 처리: 즉시·정확·개인정보 무전송.", 0, True),
    ("지식 질문('기초연금이 뭐야?') — 클라우드 LLM(Gemini 2.5 Flash, 실패 시 Groq→Claude 자동 폴백)이 실시간 답변. 화면에 '🧠 실시간 AI' 정직 라벨.", 0, True),
    ("백엔드가 없어도 전 기능 동작 — 규칙 폴백은 결함이 아니라 신뢰성 설계입니다.", 0),
    ("열면 먼저 보고 — 담아둔 복지의 마감·서류·갱신 중 급한 것부터 브리핑(능동적 개입).", 0),
])

# ════════════════════════ PART 4 다국어 ════════════════════════
section(4, "국경 없는 접근 — 7개 언어 온디바이스 AI", "외국인·다문화 사각지대를 정면으로 풉니다")

s = img_slide("온디바이스 다국어 의미검색 (헤드라인 기술)", DOC("diagram2-AI검색흐름.png"),
    caption="multilingual-e5 신경망이 브라우저 안에서 직접 실행 — 질의만 임베딩, 정책 벡터는 빌드 시 사전계산. 서버 전송 0.")

s = img_slide("영어로 검색하면", CAP("41-i18n-en.png"),
    caption="'I lost my job and need money' → 한국 복지를 뜻으로 매칭. 언어 자동감지 + AI 답변 카드(검색결과 기반 · 환각 없음)")

s = img_slide("외국어 딥퍼널 — 상세·신청까지 자국어 UI", CAP("42-i18n-en-drawer.png"),
    caption="정책 상세·신청 키트의 UI 골격이 자국어로(en·vi·zh·ja·th·ru·ar). 제도 본문은 한국어 유지 + 통역 연결(기계번역 환각 방지)")

s = img_slide("🆕 창구 도움 카드 — 말이 안 통해도 '보여주면' 됩니다", SUB("창구도움카드-영어UI.png"),
    caption="위: 본인용 자국어 안내(검증된 사전) · 아래: 직원이 읽을 큰 한국어 문장 + 필요서류 체크리스트 + 통역 연락처 · 큰 글씨 인쇄",
    note="아랍어 사용자에겐 다누리(미지원) 대신 129만 안내 — 거짓 통역 약속을 하지 않습니다")

s = base_slide("다국어 정직성 원칙 — 번역하지 않을 자유")
bullets(s, [
    ("UI 골격만 번역 — 사전에 검증해 둔 정적 문자열(8개 언어)만 사용합니다.", 0, True),
    ("제도명·금액·자격 본문은 한국어 유지 — 기계번역이 만들 수 있는 환각(오역된 자격·금액)을 원천 차단.", 0),
    ("대신 연결 — 원문 그대로 구글 번역으로 보내는 링크 + 다누리콜(13개 언어 통역 ☎1577-1366) + ☎129.", 0),
    ("음성도 자국어로 — 마이크로 묻고, 해당 언어 보이스가 있으면 음성으로 답합니다(없으면 정직하게 텍스트 안내).", 0),
    ("AI 답변 카드도 검색결과 기반 템플릿(8언어) — 생성형 환각이 구조적으로 불가능합니다.", 0, True, GREEN_D),
])

# ════════════════════════ PART 5 신청 실행 ════════════════════════
section(5, "신청 — 링크가 아니라 실행", "다른 서비스는 여기서 링크를 줍니다. 모두봄은 직접 갑니다.")

s = img_slide("경쟁 서비스와의 차이", DOC("chart4-경쟁비교.png"),
    caption="복지로/보조금24·상용 서비스는 '안내'에서 멈춥니다 — 모두봄은 민간재단·온디바이스·신청 실행·사후관리까지")

s = base_slide("신청 키트 — 어떤 브라우저·폰에서든")
bullets(s, [
    ("공식 신청 딥링크 — 복지로 신청 페이지·정부24 민원으로 직행(전 링크 실측 검증, 자동 재검증 스크립트 상시 가동).", 0),
    ("원터치 신청 — [신청] 한 번에 내 정보(이름)를 복사하고 공식 페이지를 엽니다. 간편인증만 하면 끝.", 0),
    ("신청 자동화 흐름 스테퍼 — 추천·정보작성·서류는 '자동', 인증·제출은 '본인'임을 화면에 정직하게 표시.", 0),
    ("복귀 확인 — 정부 탭에서 돌아오면 '완료하셨나요?' 1탭. 사용자가 '네'를 눌러야만 기록(클릭=완료 날조 금지).", 0, False, PEACH),
])

s = img_slide("서류 자동발급 — 실제 정부24를 조작합니다", DOC("diagram3-발급자동화.png"),
    caption="Playwright + 실제 크롬(CDP) — 로그인·메뉴 이동·양식 입력·발급까지 자동. 본인인증(카카오·PASS·네이버·토스)만 폰에서")

s = img_slide("지원 서류 15종", DOC("chart3-발급자동화.png"),
    caption="정부24 13종(등본·초본·가족관계·장애인·소득금액·납세·수급자·한부모·국세·출입국·병적·건보납부 등) + 건보 자격득실 + 고용보험 이력")

s = img_slide("서류 준비 도우미 — 웹에서 그대로", CAP("13-doc-center.png"),
    caption="필요 서류를 모아 무설치 전자증명서 발급으로 연결, '이어서 발급'이 남은 서류를 순서대로 — 발급 완료는 기기가 기억")

s = img_slide("데스크탑 앱 — '전부 자동발급' 한 번에", SUB("슬라이드8-서류도우미-데스크탑앱.png"),
    caption="한 번 인증으로 여러 서류 연쇄 발급(journey 엔진) — PDF가 바탕화면에 저장됩니다. Windows 원클릭 설치")

s = base_slide("human-in-the-loop — '안 되는 것'을 정직하게")
bullets(s, [
    ("완전 무인 자동 제출은 설계상 하지 않습니다 — 본인인증은 법이 사람에게 요구하는 단계이고, 최종 제출은 비가역·법적 행위이기 때문입니다.", 0, True),
    ("그래서 흐름의 끝은 항상 사람 — 인증 1탭, 제출 1클릭만 남기고 나머지를 전부 자동화합니다.", 0),
    ("이 경계는 한계가 아니라 신뢰입니다 — 취약계층 대상 서비스가 지켜야 할 안전장치라고 믿습니다.", 0, True, GREEN_D),
])

s = base_slide("🆕 마지막 1cm — 쓰지 못하는 분을 위한 대필")
bullets(s, [
    ("신청 사유서 — 긴급복지·장학·재단 신청의 최대 장벽 '사유 쓰기'를 대신: 입력한 상황만 문장으로(환각 0), 복사·인쇄해 바로 제출.", 0, True),
    ("위임장 — 거동이 불편하면 가족이 대신 신청하도록 위임장 초안까지.", 0),
    ("📞 전화 대본 — 129·담당기관에 전화할 때 그대로 읽는 문장 + 꼭 물어볼 체크리스트 + 메모칸.", 0, True),
    ("주민센터 방문 키트 — 정책 1건을 큰 글씨 A4로: 창구에서 할 말·서류 체크리스트·담당 전화 + 가까운 주민센터 지도.", 0),
    ("전부 규칙 기반 — 지어낸 문장이 없고, '초안이니 확인 후 사용'을 명시합니다.", 0, False, PEACH),
])

# ════════════════════════ PART 6 사후관리 ════════════════════════
section(6, "사후관리 — 신청하고 끝이 아닙니다", "서류 미비·진행 점검·갱신 임박을 에이전트가 감시합니다")

s = img_slide("나의 복지 — 상태 관리와 모니터링", CAP("12-my-full.png"),
    caption="관심→준비→신청→수급 상태 관리 + '다음 할 일' 자동 산출 + 현금성만 정직 합산한 월 최대 지원")

s = base_slide("모니터링 엔진 — 날조 없는 사후관리")
bullets(s, [
    ("사용자 기록(신청일·점검일) + 정책 갱신 주기로 산출 — 서류 미비·신청 권유·진행 점검·갱신 임박.", 0, True),
    ("정부 서버의 실시간 심사 상태는 사용자 세션 없이는 알 수 없습니다 — 그래서 '아는 척'하지 않고 공식 조회 링크로 안내합니다.", 0, False, PEACH),
    ("과거 버전의 가짜 상태 생성기(random)는 발견 즉시 제거했습니다 — 정직성은 기능보다 우선입니다.", 0),
    ("복지 캘린더 — 준비·점검·갱신 일정을 .ics로 폰 캘린더에(과거 일정은 정직하게 제외).", 0),
])

s = base_slide("가족 도움 링크 — 프라이버시를 지키는 연대")
bullets(s, [
    ("어르신의 결과를 링크 하나로 가족에게 — 받은 폰에서 같은 결과를 '재계산'합니다(서버 저장 없음, 이름 미포함).", 0, True),
    ("도우미 모드 격리 — 남의 결과를 보는 동안 내 관심목록·기록·신원이 섞이지 않도록 모든 쓰기 경로를 차단(감사로 4중 검증).", 0),
    ("🆕 글로 복사 — 결과를 카톡용 요약 글로 한 번에: 어르신이 자녀에게 '이거 봐줘' 하는 실제 경로.", 0, True),
    ("현장 모드 — 복지관 상담용 '다음 분 상담 시작'이 이전 상담자의 모든 흔적(프로필·클립보드까지)을 지웁니다.", 0),
])

# ════════════════════════ PART 7 접근성 ════════════════════════
section(7, "누구나 쓸 수 있게 — 접근성", "가장 필요한 분들이 가장 쉽게")

s = img_slide("어르신 모드 — 큰글씨·고대비", CAP("m4-elderly-on.png"),
    caption="65+ 감지 시 원탭 제안 — 큰글씨는 본문·제목까지 고르게 확대, 고대비는 히어로 강조어까지 실제로 진하게(실측 검증)")

s = base_slide("음성으로 묻고, 음성으로 듣습니다")
bullets(s, [
    ("음성 입력 — 검색·한 문장 분석·챗 어디서나 마이크로(Web Speech).", 0),
    ("읽어주기(TTS) — 결과 요약·정책 상세를 음성으로. 외국어는 해당 언어 보이스로(없으면 정직하게 안내).", 0),
    ("스크린리더 — 결과 도착·상태 변화를 라이브 영역으로 안내, 모달 포커스 트랩·ESC·aria 전면 적용.", 0),
    ("WCAG AA 대비 — 저대비 조합을 감사로 전수 색출·수정(4.5:1), reduced-motion 존중.", 0),
])

s = img_slide("PWA — 설치하면 앱, 꺼져도 동작", CAP("m1-home.png"),
    caption="홈 화면 설치·오프라인 분석(시드 190건)·자동 업데이트 — 행사장 와이파이가 끊겨도 시연이 멈추지 않습니다")

s = base_slide("프라이버시 구조 — '안 보내는' 설계")
bullets(s, [
    ("프로필·분석·검색 임베딩 — 전부 브라우저 안에서 계산하고 기기에만 저장합니다.", 0, True),
    ("주민등록번호 등 민감정보는 아예 받지 않습니다. RPA용 이름·생년월일도 기기에만.", 0),
    ("선택 동기화(카카오·구글 로그인)를 켜도 '나의 복지' 목록만 본인 계정 클라우드로 — 언제든 삭제.", 0),
    ("취약계층이 안심하고 쓰는 것 — 그 자체가 접근성이라고 생각합니다.", 0, True, GREEN_D),
])

# ════════════════════════ PART 8 품질·정직성 ════════════════════════
section(8, "품질·정직성 그리고 실용성", "과장은 기능이 아니라 버그로 취급했습니다 — 그리고 오늘 바로 쓸 수 있습니다")

s = base_slide("검증 수치 (전부 실측)")
stat_cards(s, [("809", "자동 테스트\n(프론트 679 · 백엔드 130)"), ("13여정", "실브라우저 E2E\n(모바일·도우미·다국어 포함)"), ("17라운드", "다중에이전트 적대 감사"), ("90여 건", "확정 결함 발견·정정")], y=1.7, ch=2.0, size_v=30)
bullets(s, [
    ("적대 감사 = AI 에이전트 수십 개가 '결함을 찾도록' 경쟁하고, 별도 에이전트가 재현 가능성을 검증 — 확정된 것만 수정했습니다.", 0),
    ("모든 수정은 회귀 테스트로 고정 — 같은 결함이 되돌아오지 않습니다.", 0),
    ("헤드라인 금액 9종은 2026 공식 출처로 재검증(기초연금 349,700원·생계급여 중위32% 등 전건 일치), 공식 딥링크 19종 자동 생존 점검.", 0),
], y=4.0, size=14)

s = base_slide("정직성을 코드로 — 대표 장치들")
bullets(s, [
    ("현금성 게이트 — 서비스 한도(장기요양 월 251만)·바우처·감면·대출·임금대체를 '현금'으로 합산·표기하지 않습니다. 위반 시 테스트 실패.", 0, True),
    ("심사·선발형(민간재단)은 '자격 있음' 단정 금지 — 표기 자체를 테스트가 거부합니다.", 0),
    ("모르면 모른다고 — 소득을 안 적으면 '중위 80% 확정'이 아니라 '소득 확인 필요' 조건부로 안내합니다.", 0),
    ("클릭≠완료 — 신청·발급은 사용자가 '네'를 눌러야만 기록됩니다.", 0),
    ("결과 인쇄물·공유 텍스트·시뮬레이터까지 — 모든 표시면이 같은 게이트를 공유합니다.", 0),
])

s = img_slide("데이터 파이프라인 — 지속 가능한 최신성", DOC("diagram6-데이터파이프라인.png"),
    caption="한국사회보장정보원 공공데이터(OpenAPI) → 정규화 ETL → 런타임 병합. 재수집 한 번으로 전체 갱신")

s = base_slide("지금 바로 사용 가능")
bullets(s, [
    ("라이브 서비스 운영 중 — biocode67.github.io/modoo-bom (설치·로그인·비용 없음).", 0, True),
    ("운영비 ≈ 0원 — 정적 호스팅 + 온디바이스 구조라 사용자가 늘어도 비용이 늘지 않습니다.", 0),
    ("Windows 데스크탑 앱 원클릭 설치(자동발급), 크롬 확장, 서버사이드 체험 모드까지 3채널.", 0),
    ("복지관 현장 모드 — 상담사가 어르신을 연속 상담할 수 있는 초기화·인쇄·전화 대본 도구 내장.", 0),
])

s = base_slide("확장 로드맵 — 그리고 지켜온 원칙")
bullets(s, [
    ("보급 — 지자체·복지관·다문화가족지원센터에 무상 보급(무운영비라 즉시 가능).", 0),
    ("정부 연계 — 전자증명서(전자문서지갑)·마이데이터 연동으로 '서류 없는 신청'까지. 같은 에이전트 루프는 세금 감면·교육 지원 등 '몰라서 놓치는 권리' 전반으로 확장됩니다.", 0),
    ("원칙 ① 실데이터만 — 가짜 정책·금액·상태를 만들지 않습니다.", 0, True),
    ("원칙 ② 실제 동작만 — 목업이 아니라 실제 정부24·실제 신경망·실제 배포로 검증합니다.", 0, True),
    ("원칙 ③ 사람의 자리 — 인증·제출·최종 확인은 언제나 사용자의 것입니다.", 0, True),
    ("이 문서의 모든 수치는 라이브 화면·코드·테스트로 재현 가능합니다.", 0, False, PEACH),
])

# ════════════════════════ 마무리 ════════════════════════
s = prs.slides.add_slide(BLANK); page_no += 1
bg = s.shapes.add_shape(1, 0, 0, W, H); _fill(bg, GREEN_D)
tb = s.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.4))
tf = tb.text_frame
_text(tf, [[("몰라서 못 받는 복지가 없는 나라를 향해", False, RGBColor(0xBB, 0xF7, 0xD0))]], 18, WHITE, space_after=14)
p = tf.add_paragraph(); r = p.add_run(); r.text = "모두봄 — 모두의 봄날을 위한 복지 도우미 🌱"
r.font.name = FONT; r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = WHITE
p = tf.add_paragraph(); p.space_before = Pt(22)
r = p.add_run(); r.text = "지금 체험: biocode67.github.io/modoo-bom"
r.font.name = FONT; r.font.size = Pt(20); r.font.color.rgb = WHITE
p = tf.add_paragraph(); p.space_before = Pt(8)
r = p.add_run(); r.text = "소스: github.com/BioCode67/modoo-bom  ·  본 산출물의 전 수치는 실측·재현 가능"
r.font.name = FONT; r.font.size = Pt(13); r.font.color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
ft = s.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.35))
_text(ft.text_frame, f"모두봄        {page_no:02d} / 50", 10, RGBColor(0xD1, 0xFA, 0xE5))

OUT = A("docs", "제출", "모두봄-최종제출-50p.pptx")
prs.save(OUT)
print(f"슬라이드 {page_no}장 생성 → {OUT}")
