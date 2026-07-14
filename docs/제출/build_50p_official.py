# -*- coding: utf-8 -*-
"""모두봄 — 공식 양식 기반 50p 최종제출용 생성기.
베이스: 모두봄-예선산출물-작성본.pptx(공식 양식 + 팀 정보 완성본).
방식: 양식 프레임(상하단 라인·섹션 라벨·챕터 태그·소제목·로고·쪽번호)을 그대로 복제하고
본문만 교체 — 양식 임의 변경 없음. 콘텐츠 수치는 전부 실측.
실행: python docs/제출/build_50p_official.py
"""
import sys, os, copy
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
A = lambda *p: os.path.join(ROOT, *p)
BASE_FILE = A("docs", "제출", "모두봄-예선산출물-작성본.pptx")
OUT = A("docs", "제출", "모두봄-최종제출-50p.pptx")

GREEN_D = RGBColor(0x15, 0x80, 0x3D)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
INK = RGBColor(0x20, 0x29, 0x24)
GRAY = RGBColor(0x5A, 0x64, 0x5D)
PEACH = RGBColor(0xC2, 0x41, 0x0C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"
IN = lambda v: Inches(v)

prs = Presentation(BASE_FILE)
W, H = prs.slide_width, prs.slide_height

# ── 슬라이드 복제(이미지 rel 재바인딩 포함) ─────────────────────────
def clone_slide(src_idx):
    src = prs.slides[src_idx]
    new = prs.slides.add_slide(src.slide_layout)
    # add_slide가 만든 기본 placeholder 제거(원본 것을 그대로 복제)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    # 이미지 rel 매핑 준비
    img_map = {}
    for rid, rel in src.part.rels.items():
        if "image" in rel.reltype:
            new_rid = new.part.relate_to(rel.target_part, rel.reltype)
            img_map[rid] = new_rid
    for sh in src.shapes:
        el = copy.deepcopy(sh._element)
        # blip r:embed 재바인딩
        for blip in el.iter(qn('a:blip')):
            old = blip.get(qn('r:embed'))
            if old in img_map:
                blip.set(qn('r:embed'), img_map[old])
        new.shapes._spTree.append(el)
    return new

# ── 프레임만 남기고 본문 비우기 + 소제목 교체 ───────────────────────
def strip_body(slide, subtitle):
    """유지: 라인·로고(하단)·쪽번호·상단 라벨/태그(top<1.0)·본문 테두리·소제목 '단 하나'.
       삭제: 그 외 본문 콘텐츠 전부. 소제목은 문단을 XML로 싹 비우고 한 줄로 재구성(겹침 방지)."""
    subtitle_done = False
    for sh in list(slide.shapes):
        t = Emu(sh.top).inches if sh.top is not None else 0
        l = Emu(sh.left).inches if sh.left is not None else 0
        w = Emu(sh.width).inches if sh.width is not None else 0
        h = Emu(sh.height).inches if sh.height is not None else 0
        st = str(sh.shape_type)
        keep = False
        if 'LINE' in st and t < 1.9: keep = True           # 상단 라인·소제목 밑줄(하단 라인은 t≈7.0 → 아래 조건)
        elif 'LINE' in st and t > 6.5: keep = True         # 하단 라인
        elif 'PLACEHOLDER' in st: keep = True              # 쪽번호
        elif 'PICTURE' in st and t > 6.5: keep = True      # 우하단 로고
        elif t < 1.0: keep = True                          # 상단 라벨·챕터 태그
        elif (not subtitle_done and sh.has_text_frame
              and abs(l - 1.0) < 0.35 and abs(t - 1.3) < 0.35 and h < 0.5):
            # 소제목(양식 위치의 단 하나) — 문단 전부 제거 후 한 줄 재구성 + 폭 확장(줄바꿈 겹침 방지)
            keep = True; subtitle_done = True
            tf = sh.text_frame
            base_run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else None
            sz = base_run.font.size if base_run is not None and base_run.font.size else Pt(16)
            for p in list(tf.paragraphs[1:]):
                p._p.getparent().remove(p._p)
            p0 = tf.paragraphs[0]
            for r in list(p0.runs):
                r._r.getparent().remove(r._r)
            r = p0.add_run(); r.text = subtitle
            r.font.name = FONT; r.font.size = sz; r.font.bold = True; r.font.color.rgb = INK
            tf.word_wrap = False
            sh.width = IN(11.5)
        elif abs(l - 0.6) < 0.2 and abs(t - 1.1) < 0.2 and w > 11.0: keep = True  # 본문 테두리 프레임
        if not keep:
            sh._element.getparent().remove(sh._element)
    return slide

# ── 콘텐츠 헬퍼(본문 영역: x 0.9~12.4 / y 1.95~6.35) ────────────────
def add_bullets(slide, items, x=1.0, y=2.0, w=11.4, h=4.3, size=14, gap=8):
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        text = it[0] if isinstance(it, tuple) else it
        lvl = it[1] if isinstance(it, tuple) and len(it) > 1 else 0
        bold = it[2] if isinstance(it, tuple) and len(it) > 2 else False
        col = it[3] if isinstance(it, tuple) and len(it) > 3 else INK
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = lvl
        r = p.add_run()
        r.text = ('· ' if lvl == 0 else '– ') + text
        r.font.name = FONT; r.font.size = Pt(size if lvl == 0 else size - 1.5)
        r.font.bold = bold; r.font.color.rgb = col
    return tb

def add_img(slide, img, caption=None, note=None, max_h=4.15):
    from PIL import Image
    if not os.path.exists(img):
        tb = slide.shapes.add_textbox(IN(1.0), IN(3.2), IN(11.4), IN(0.6))
        tf = tb.text_frame; p = tf.paragraphs[0]; r = p.add_run()
        r.text = f"[이미지: {os.path.basename(img)}]"; r.font.name = FONT; r.font.size = Pt(12); r.font.color.rgb = GRAY
        return
    iw, ih = Image.open(img).size
    maxw = IN(11.4); maxh = IN(max_h if not (caption and note) else max_h - 0.3)
    ratio = min(maxw / iw, maxh / ih)
    pw, ph = int(iw * ratio), int(ih * ratio)
    left = int((W - pw) / 2)
    pic = slide.shapes.add_picture(img, left, IN(1.95), width=pw, height=ph)
    # 스크린샷 경계 명확화 — 연회색 얇은 테두리(디자인 개선)
    pic.line.fill.solid(); pic.line.fill.fore_color.rgb = RGBColor(0xC8, 0xD2, 0xCB); pic.line.width = Pt(0.75)
    ytxt = 1.95 + Emu(ph).inches + 0.12
    if caption:
        tb = slide.shapes.add_textbox(IN(0.9), IN(ytxt), IN(11.6), IN(0.4))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption
        r.font.name = FONT; r.font.size = Pt(12); r.font.color.rgb = GRAY
        ytxt += 0.38
    if note:
        tb = slide.shapes.add_textbox(IN(0.9), IN(ytxt), IN(11.6), IN(0.4))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = note
        r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GREEN_D

def add_stats(slide, cards, y=2.2, ch=1.7, size_v=26, size_l=12):
    n = len(cards)
    gap = IN(0.25)
    total_w = IN(11.4)
    cw = int((total_w - gap * (n - 1)) / n)
    x = IN(1.0)
    for v, l in cards:
        box = slide.shapes.add_shape(1, x, IN(y), cw, IN(ch))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.fill.solid(); box.line.fill.fore_color.rgb = GREEN; box.line.width = Pt(1.2)
        tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = v; r.font.name = FONT; r.font.size = Pt(size_v); r.font.bold = True; r.font.color.rgb = GREEN_D
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = l; r2.font.name = FONT; r2.font.size = Pt(size_l); r2.font.color.rgb = GRAY
        x += cw + gap

CAP = lambda *p: A("frontend", "e2e", "audit", *p)
DOC = lambda *p: A("docs", "기획서자료", *p)
SUB = lambda *p: A("docs", "제출", "스크린샷", *p)

# ── 표준 캔버스: 작성본 6번(최종 산출물 개요, idx 5)을 복제해 사용 ──
STD = 5  # 0-based: 슬라이드 6

def content_slide(subtitle):
    s = clone_slide(STD)
    strip_body(s, subtitle)
    return s

def bullet_slide(subtitle, items, **kw):
    s = content_slide(subtitle)
    add_bullets(s, items, **kw)
    return s

def image_slide(subtitle, img, caption=None, note=None, max_h=3.9):
    s = content_slide(subtitle)
    add_img(s, img, caption, note, max_h)
    return s

# ═══════════════ 본문 43장(최종 산출물 확장) ═══════════════
# ── A. 최종 산출물 개요 (6장) ──
s = bullet_slide("최종 산출물 개요 ① 문제 — 몰라서 못 받는 복지", [
    ("자격이 있어도 놓치는 3가지 이유", 0, True),
    ("① 몰라서 — 중앙부처 460 + 지자체 4,683 + 민간재단·서민금융… 한 곳에서 볼 수 없습니다.", 0),
    ("② 어려워서 — '소득인정액'·'중위소득 48%'·'차상위' 같은 행정 용어의 벽.", 0),
    ("③ 못 써서 — 신청서의 '사유'를 못 쓰고, 창구에서 말이 막히고, 서류에서 포기합니다.", 0),
    ("가장 필요한 분들(어르신·장애인·외국인·위기가구)이 가장 큰 장벽을 만납니다.", 0, True, PEACH),
    ("모두봄 = 찾기→이해→행동(신청)→사후관리 전 여정을 대행하는 개인 복지 AI Agent입니다.", 0, True, GREEN_D),
])

s = bullet_slide("최종 산출물 개요 ② 한 사람의 이야기 — 김복순 님(72세, 독거)", [
    ("기초연금은 받고 계십니다. 그런데 그게 전부일까요?", 0, True),
    ("'72세 혼자 사는데 소득이 적어요' 한 문장이면 —", 0),
    ("핵심 현금지원 월 최대 64만원(중복수급 미반영), 강력 추천 10여 건: 주거급여·의료급여·틀니 지원·노인 일자리·에너지바우처…", 1, True),
    ("정부 데이터 어디에도 없는 민간재단(심장 수술비·위기가정 지원)까지 함께 찾아드립니다.", 1),
    ("김복순 님 같은 분이 전국에 수백만 명 — 모두봄은 그분들의 '몰랐던 권리'를 찾아드립니다.", 0, True, GREEN_D),
])

s = image_slide("최종 산출물 개요 ③ 경쟁 서비스와의 차이", DOC("chart4-경쟁비교.png"),
    caption="복지로/보조금24·상용 서비스는 '안내'에서 멈춥니다 — 모두봄은 민간재단·온디바이스·신청 실행·사후관리까지")

s = content_slide("최종 산출물 개요 ④ 검증 수치 (전부 실측)")
add_stats(s, [("5,300+", "실데이터 복지(건)"), ("809", "자동 테스트(679+130)"), ("15종", "서류 자동발급"), ("7개 언어", "온디바이스 AI 검색")], y=2.1, ch=1.6)
add_bullets(s, [
    ("실브라우저 E2E 13여정(모바일·도우미·다국어 포함) · 다중에이전트 적대 감사 17라운드 · 확정 결함 90여 건 정정.", 0),
    ("헤드라인 금액 9종은 2026 공식 출처로 재검증(기초연금 349,700원·생계급여 중위32% 등 전건 일치).", 0),
    ("공식 딥링크 19종 자동 생존 점검 — 본 문서의 모든 수치는 라이브·코드·테스트로 재현 가능합니다.", 0, True, GREEN_D),
], y=4.1, size=13)

s = bullet_slide("최종 산출물 개요 ⑤ 정직성을 코드로", [
    ("과장은 기능이 아니라 버그로 취급했습니다 — 정직성 장치가 테스트로 강제됩니다.", 0, True),
    ("현금성 게이트 — 서비스 한도(장기요양 월 251만)·바우처·감면·대출·임금대체를 '현금'으로 합산·표기하지 않습니다.", 0),
    ("심사·선발형(민간재단)은 '자격 있음' 단정 금지 — 표기 자체를 테스트가 거부합니다.", 0),
    ("모르면 모른다고 — 소득 미입력 시 '중위 80% 확정'이 아니라 '소득 확인 필요' 조건부 안내.", 0),
    ("클릭≠완료 — 신청·발급은 사용자가 '네'를 눌러야만 기록됩니다.", 0),
    ("AI 답변도 검색결과 기반 템플릿 — 생성형 환각이 구조적으로 불가능합니다.", 0),
])

# ── B. AI Agent 작동 구조 (7장) ──
s = bullet_slide("AI Agent 작동 구조 ① 챗봇이 아니라 '루프'", [
    ("인지 — 프로필·자연어·음성에서 상황 신호를 읽습니다(실직·한부모·장애·위기 등 사각지대 신호 포함).", 0, True),
    ("판단 — 2026 정밀 선정기준으로 자격 판정, 우선순위·예상 금액·'왜 나에게 맞는지' 근거 생성.", 0, True),
    ("행동 — 한 번에 담기, 서류 자동발급(RPA), 신청서 정보 자동 입력, 사유서·위임장·전화 대본 대필.", 0, True),
    ("관찰 — 신청 진행·서류 미비·갱신 임박·마감을 감시하고 먼저 보고(챗을 열면 브리핑부터).", 0, True),
    ("인증·최종 제출만 사람에게 — human-in-the-loop 설계 원칙.", 0, False, PEACH),
])

s = image_slide("AI Agent 작동 구조 ② 전체 아키텍처 — 3층 지능", DOC("diagram1-아키텍처.png"),
    caption="① 온디바이스 신경망(다국어 의미검색·프라이버시) ② LangGraph 10노드 에이전트 ③ Playwright RPA(실제 정부24)")

s = image_slide("AI Agent 작동 구조 ③ LangGraph 10노드", DOC("diagram5-LangGraph10노드.png"),
    caption="프로필 분석 → RAG 검색 → 자격 판별 → 재판정 루프 → 가이드·서류·포트폴리오·알림·추적 → 오케스트레이터",
    note="라이브에서 실시간 스트리밍으로 10노드가 순차 점등 — 연출이 아니라 실제 실행입니다")

s = image_slide("AI Agent 작동 구조 ④ 데이터 — 약 5,300건 실데이터", DOC("chart1-출처구성.png"),
    caption="정부 시드 124(정밀 규칙·검증 금액) + 지원사업 33 + 주택 7 + 서민금융 5 + 민간재단 21 + 공공데이터(중앙 460·지자체 4,683)",
    note="가짜 데이터 미생성 원칙 — 전 항목 실데이터")

s = image_slide("AI Agent 작동 구조 ⑤ 데이터 파이프라인", DOC("diagram6-데이터파이프라인.png"),
    caption="한국사회보장정보원 공공데이터(OpenAPI) → 정규화 ETL → 런타임 병합 — 재수집 한 번으로 전체 갱신")

s = bullet_slide("AI Agent 작동 구조 ⑥ 민간재단 큐레이션 21건", [
    ("복지로·보조금24 어디에도 없는 기업·재단 지원을 하나하나 실측 검증해 수록했습니다.", 0, True),
    ("장학: 현대차 정몽구 스칼러십 · 관정 · 삼성꿈장학 · 미래에셋", 0),
    ("의료·위기: 한국심장재단 · 백혈병어린이재단 · 아산 SOS · 초록우산 · 밀알(장애·재활)", 0),
    ("정직한 표시 — 심사·선발형임을 명시하고 '자격 있음'으로 단정하지 않습니다(테스트 강제).", 0, False, PEACH),
    ("정책서민금융 5건(햇살론 등)도 '대출'임을 명시하고 현금 합산에서 제외합니다.", 0),
])

s = bullet_slide("AI Agent 작동 구조 ⑦ 챗 하이브리드 — 행동은 로컬, 지식은 LLM", [
    ("행동 의도(담기·서류·자격·브리핑) — 기기 안 규칙 에이전트: 즉시·정확·개인정보 무전송.", 0, True),
    ("지식 질문('기초연금이 뭐야?') — 클라우드 LLM(Gemini 2.5 Flash, 실패 시 Groq→Claude 자동 폴백). 화면에 '🧠 실시간 AI' 정직 라벨.", 0, True),
    ("백엔드가 없어도 전 기능 동작 — 규칙 폴백은 결함이 아니라 신뢰성 설계입니다.", 0),
    ("열면 먼저 보고 — 담아둔 복지의 마감·서류·갱신 중 급한 것부터 브리핑(능동적 개입).", 0),
])

# ── C. 핵심 기능 설명 (18장) ──
s = image_slide("핵심 기능 ① 대화형 온보딩 — '새싹이와 대화'", CAP("m2-analyze-chat.png"),
    caption="폼 대신 마스코트가 하나씩 묻고 탭으로만 답합니다. 어르신(65+)에겐 큰글씨 원탭 제안, 새로고침해도 이어서")

s = image_slide("핵심 기능 ② 한 문장 분석 (음성 지원)", SUB("한문장분석결과-2026-07-15.png"),
    caption="'72세 혼자 사는데 소득이 적어요' → 핵심 현금지원 월 최대 64만원 — 2026-07-15 라이브 실측 화면 그대로(무보정)")

s = image_slide("핵심 기능 ③ 결과 — 에이전트가 '다음 행동'을 제안", CAP("10-result-top.png"),
    caption="'이렇게 이해했어요' 신호 투명화 + 추천 한 번에 담기 + 마감 임박·아깝게 놓친 복지·연계 감면까지")

s = image_slide("핵심 기능 ④ 정책 상세 — '내가 받을 수 있는 이유'", CAP("d1-detail-top.png"),
    caption="조건 나열이 아니라 내 정보와 대조한 체크리스트(만 72세 ✓ · 소득 하위 25% ✓)로 근거 제시")

s = bullet_slide("핵심 기능 ⑤ 숨은 지능 — 추천에서 끝나지 않는 판단", [
    ("수급 조합 도우미 — '하나만'(기초연금↔장애인연금) · '감액 주의'(기초→생계) · '병급 가능'(아동수당+부모급여). 공식 규정 실측.", 0),
    ("아깝게 놓친 복지 — 딱 한 조건만 벗어난 정책을 근거와 함께(판단 근거 노출).", 0),
    ("연계 감면 일괄 안내 — 자격이 되면 딸려오는 통신·전기·가스·TV수신료 감면.", 0),
    ("생애 타임라인·시뮬레이터 — '만 65세가 되면'·'아이가 태어나면' 무엇이 열리고 닫히는지 미리.", 0),
    ("긴급복지 빠른 진단 — 위기 상황이면 129·1366·1393 핫라인 최우선(안전 우선).", 0, False, PEACH),
])

s = image_slide("핵심 기능 ⑥ 행동형 챗 — '다 담아줘' 한마디", CAP("32-chat-result.png"),
    caption="대화 맥락을 기억해 방금 보여준 복지를 저장 — 행동은 기기 안 에이전트가 즉시 수행")

s = image_slide("핵심 기능 ⑦ 온디바이스 다국어 AI 의미검색 (헤드라인)", DOC("diagram2-AI검색흐름.png"),
    caption="multilingual-e5 신경망이 브라우저에서 직접 실행 — 질의만 임베딩, 서버 전송 0")

s = image_slide("핵심 기능 ⑧ 영어로 검색하면", CAP("41-i18n-en.png"),
    caption="'I lost my job and need money' → 한국 복지를 '뜻'으로 매칭 + AI 답변 카드(검색결과 기반·환각 없음)")

s = image_slide("핵심 기능 ⑨ 외국어 딥퍼널 — 상세·신청까지 자국어 UI", CAP("42-i18n-en-drawer.png"),
    caption="정책 상세·신청 키트 UI 골격이 자국어로(en·vi·zh·ja·th·ru·ar). 제도 본문은 한국어 유지 + 통역 연결(환각 방지)")

s = image_slide("핵심 기능 ⑩ 창구 도움 카드 — 말이 안 통해도 '보여주면' 됩니다", SUB("창구도움카드-영어UI.png"),
    caption="본인용 자국어 안내(검증된 사전) + 직원이 읽을 큰 한국어 + 서류 체크리스트 + 통역 연락처 · 큰 글씨 인쇄",
    note="아랍어 사용자에겐 다누리(미지원) 대신 129만 안내 — 거짓 통역 약속을 하지 않습니다", max_h=3.9)

s = bullet_slide("핵심 기능 ⑪ 신청 키트 — 어떤 브라우저·폰에서든", [
    ("공식 신청 딥링크 — 복지로 신청 페이지·정부24 민원 직행(전 링크 실측 검증·자동 재검증 상시).", 0),
    ("원터치 신청 — [신청] 한 번에 내 정보(이름) 복사 + 공식 페이지 열기. 간편인증만 하면 끝.", 0),
    ("신청 자동화 스테퍼 — 추천·정보작성·서류는 '자동', 인증·제출은 '본인'임을 화면에 정직하게 표시.", 0),
    ("복귀 확인 — 정부 탭에서 돌아오면 '완료하셨나요?' 1탭. '네'를 눌러야만 기록(날조 금지).", 0, False, PEACH),
])

s = image_slide("핵심 기능 ⑫ 서류 자동발급 — 실제 정부24를 조작합니다", DOC("diagram3-발급자동화.png"),
    caption="Playwright + 실제 크롬(CDP) — 로그인·양식 입력·발급까지 자동. 본인인증(카카오·PASS·네이버·토스)만 폰에서")

s = image_slide("핵심 기능 ⑬ 지원 서류 15종", DOC("chart3-발급자동화.png"),
    caption="정부24 13종(등본·초본·가족관계·장애인·소득금액·납세·수급자·한부모·국세·출입국·병적·건보납부 등) + 건보 자격득실 + 고용보험 이력")

s = image_slide("핵심 기능 ⑭ 서류 준비 도우미 — 웹에서 그대로", CAP("13-doc-center.png"),
    caption="무설치 전자증명서 발급 연결 + '이어서 발급'이 남은 서류를 순서대로 — 발급 완료는 기기가 기억")

s = image_slide("핵심 기능 ⑮ 데스크탑 앱 — '전부 자동발급' 한 번에", SUB("슬라이드8-서류도우미-데스크탑앱.png"),
    caption="한 번 인증으로 여러 서류 연쇄 발급(journey 엔진) — PDF가 바탕화면에 저장. Windows 원클릭 설치")

s = bullet_slide("핵심 기능 ⑯ 마지막 1cm — 쓰지 못하는 분을 위한 대필", [
    ("신청 사유서 — 최대 장벽 '사유 쓰기'를 대신: 입력한 상황만 문장으로(환각 0), 복사·인쇄 즉시 제출.", 0, True),
    ("위임장 — 거동이 불편하면 가족이 대신 신청하도록 초안까지.", 0),
    ("전화 대본 — 129·기관 통화가 두려운 분이 그대로 읽는 문장 + 질문 체크리스트 + 메모칸.", 0, True),
    ("주민센터 방문 키트 — 큰 글씨 A4 인쇄(창구 멘트·서류·담당 전화) + 가까운 주민센터 지도.", 0),
    ("전부 규칙 기반 — 지어낸 문장이 없고 '초안이니 확인 후 사용'을 명시합니다.", 0, False, PEACH),
])

s = image_slide("핵심 기능 ⑰ 사후관리 — 나의 복지", CAP("12-my-full.png"),
    caption="관심→준비→신청→수급 상태 관리 + '다음 할 일' 자동 산출 + 현금성만 정직 합산")

s = bullet_slide("핵심 기능 ⑱ 모니터링·가족 연대 — 날조 없는 사후관리", [
    ("사용자 기록(신청일·점검일)+정책 갱신 주기로 산출 — 서류 미비·신청 권유·진행 점검·갱신 임박.", 0, True),
    ("정부 심사 상태는 세션 없이 알 수 없습니다 — '아는 척' 대신 공식 조회 링크로 안내(가짜 상태 생성기는 발견 즉시 제거).", 0, False, PEACH),
    ("복지 캘린더 — 준비·점검·갱신 일정을 .ics로 폰 캘린더에.", 0),
    ("가족 도움 링크 — 결과를 링크로 공유(이름 미포함), 받은 폰에서 재계산. 도우미 모드는 내 데이터와 완전 격리.", 0),
    ("글로 복사 — 결과를 카톡용 요약으로 한 번에: 어르신이 자녀에게 '이거 봐줘' 하는 실제 경로.", 0, True),
])

# ── D. 사용자 이용 흐름 (4장) ──
s = image_slide("사용자 이용 흐름 ① 전체 여정", DOC("diagram4-사용자여정.png"),
    caption="발견 → 이해 → 준비 → 신청 → 사후관리 — 끊기는 지점 없이 하나의 흐름으로", max_h=3.9)

s = bullet_slide("사용자 이용 흐름 ② 시작 — 3가지 입력", [
    ("① 새싹이와 대화(기본) — 탭으로만 답하는 대화형. 어르신 큰글씨 원탭 제안.", 0),
    ("② 한 문장 — '퇴사하고 일자리 찾는 청년이에요' (음성 지원, 7개 언어).", 0),
    ("③ 직접 입력 — 1분 위저드(페르소나 원탭 시작 지원).", 0),
    ("어떤 길로 시작해도 같은 정밀 엔진 — 입력 방식이 결과 품질을 가르지 않습니다.", 0, True, GREEN_D),
])

s = image_slide("사용자 이용 흐름 ③ 어르신 모드", CAP("m4-elderly-on.png"),
    caption="큰글씨는 본문·제목까지 고르게 확대, 고대비는 강조어까지 실제로 진하게(실측 검증) + 음성 입력·읽어주기(TTS)")

s = image_slide("사용자 이용 흐름 ④ PWA — 설치하면 앱, 꺼져도 동작", CAP("m1-home.png"),
    caption="홈 화면 설치·오프라인 분석·자동 업데이트 — 현장 와이파이가 끊겨도 이용이 멈추지 않습니다")

# ── E. 실용성 및 확장 가능성 (5장) ──
s = bullet_slide("실용성 ① 지금 바로 사용 가능", [
    ("라이브 서비스 운영 중 — biocode67.github.io/modoo-bom (설치·로그인·비용 없음).", 0, True),
    ("운영비 ≈ 0원 — 정적 호스팅 + 온디바이스 구조라 사용자가 늘어도 비용이 늘지 않습니다.", 0),
    ("Windows 데스크탑 앱 원클릭 설치 · 크롬 확장 · 서버사이드 체험 모드 — 3채널.", 0),
    ("복지관 현장 모드 — '다음 분 상담 시작'이 이전 상담자의 모든 흔적(클립보드까지)을 삭제.", 0),
])

s = bullet_slide("실용성 ② 보안·프라이버시 — '안 보내는' 설계", [
    ("프로필·분석·검색 임베딩 전부 브라우저 안에서 계산·저장(서버 전송 0).", 0, True),
    ("주민등록번호 등 민감정보는 아예 받지 않습니다.", 0),
    ("선택 동기화(카카오·구글)를 켜도 '나의 복지' 목록만 본인 계정으로 — 언제든 삭제.", 0),
    ("완전 무인 제출은 설계상 배제 — 본인인증·최종 제출은 법이 사람에게 요구하는 단계이며, 그 경계가 신뢰입니다.", 0, True, PEACH),
])

s = bullet_slide("실용성 ③ 품질 체계", [
    ("자동 테스트 809개(프론트 vitest 679 · 백엔드 pytest 130) — 모든 수정은 회귀 테스트로 고정.", 0, True),
    ("실브라우저 E2E 13여정 — 모바일·도우미·다국어·창구 카드 여정까지 매 배포 검증.", 0),
    ("다중에이전트 적대 감사 17라운드 — AI 에이전트 수십 개가 결함을 찾고 별도 에이전트가 재현성 검증, 확정 결함 90여 건 정정.", 0),
    ("데이터 정확성 — 헤드라인 금액 9종 2026 공식 출처 재검증(전건 일치), 공식 딥링크 19종 자동 생존 점검.", 0),
])

s = bullet_slide("실용성 ④ 확장 로드맵", [
    ("보급 — 지자체·복지관·다문화가족지원센터에 무상 보급(무운영비라 즉시 가능).", 0),
    ("데이터 — 공공데이터 재수집 자동화 + 민간재단 제휴로 커버리지 확대.", 0),
    ("정부 연계 — 전자증명서(전자문서지갑)·마이데이터 연동으로 '서류 없는 신청'까지.", 0),
    ("같은 에이전트 루프는 세금 감면·교육 지원·재난 지원 등 '몰라서 놓치는 권리' 전반으로 확장 가능합니다.", 0, True, GREEN_D),
])

s = bullet_slide("실용성 ⑤ 팀이 지켜온 원칙", [
    ("실데이터만 — 가짜 정책·가짜 금액·가짜 상태를 만들지 않습니다.", 0, True),
    ("실제 동작만 — 목업이 아니라 실제 정부24·실제 신경망·실제 배포로 검증합니다.", 0, True),
    ("사람의 자리 — 인증과 제출, 최종 확인은 언제나 사용자의 것입니다.", 0, True),
    ("이 문서의 모든 수치는 라이브 화면·코드·테스트로 재현 가능합니다.", 0, False, PEACH),
])

# ═══════════ 슬라이드 순서 재배치: 새 본문 43장을 10번(실용성) 뒤·부록(11) 앞으로 ═══════════
xml_slides = prs.slides._sldIdLst
slides = list(xml_slides)
# 현재: [0..10]=원본 11장, [11..]=새 43장. 목표: 원본 1~10, 새 43장, 부록(원본 11)
appendix = slides[10]
new_ones = slides[11:]
xml_slides.remove(appendix)
for el in new_ones:
    xml_slides.remove(el)
for el in new_ones:
    xml_slides.append(el)
xml_slides.append(appendix)

total = len(list(prs.slides))
prs.save(OUT)
print(f"슬라이드 {total}장 생성(양식 프레임 유지) → {OUT}")
