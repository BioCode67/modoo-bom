# -*- coding: utf-8 -*-
"""제출용 PPTX 생성 — 모두봄_포스터.pdf → 200DPI 래스터 → 공식 양식(70x140cm) 1장에 전면 배치.

왜 이미지로 넣는가:
  주최측이 "제출한 파일을 그대로 출력하여 전시"하므로, 심사자 PC에 우리 폰트가 없을 때
  글자가 치환돼 레이아웃이 깨지는 사고를 원천 차단한다. PDF(벡터)와 PPTX(래스터)가
  화면상 완전히 동일해지므로 "①과 동일한 파일을 PDF로 변환" 요건도 그대로 충족한다.

전제: 같은 폴더에 공식 양식(template.pptx)이 있어야 한다.
      (협의회 배포본 'SW부문 포스터세션 양식' — 저작권 자료라 저장소에는 포함하지 않음)
사용법: cd docs/제출/포스터 && python render.py && python build-pptx.py
"""
import pathlib
import subprocess
import sys

HERE = pathlib.Path(__file__).resolve().parent
PDF = HERE / "모두봄_포스터.pdf"
TEMPLATE = HERE / "template.pptx"
OUT = HERE / "모두봄_포스터.pptx"
RASTER = HERE / "_print-1.jpg"


def main() -> int:
    if not PDF.exists():
        print("먼저 render.py로 PDF를 만들어 주세요."); return 1
    if not TEMPLATE.exists():
        print("공식 양식(template.pptx)이 필요합니다 — 협의회 배포 파일을 이 폴더에 두세요."); return 1

    # 1) PDF → 200DPI JPEG (70x140cm 기준 5514x11026px)
    subprocess.run(["pdftoppm", "-jpeg", "-jpegopt", "quality=94", "-r", "200",
                    str(PDF), str(RASTER.with_suffix("")).replace("-1", "")],
                   check=True, stderr=subprocess.DEVNULL)

    from pptx import Presentation
    prs = Presentation(str(TEMPLATE))
    W, H = prs.slide_width, prs.slide_height

    # 2) 양식 1페이지(제출 안내)는 삭제 — 제출은 총 1장
    lst = prs.slides._sldIdLst
    first = list(lst)[0]
    prs.part.drop_rel(first.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"))
    lst.remove(first)

    # 3) 남은 슬라이드의 안내 텍스트 박스 제거 후 포스터를 전면 배치
    slide = prs.slides[0]
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    slide.shapes.add_picture(str(RASTER), 0, 0, width=W, height=H)

    prs.save(str(OUT))
    RASTER.unlink(missing_ok=True)
    print(f"완료: {OUT.name} ({OUT.stat().st_size/1e6:.1f} MB · 슬라이드 1장 · 70x140cm)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
